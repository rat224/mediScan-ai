# app.py — MediScan.AI v4 (Hospital report + Multi-language translation)
import io
import os
import re
import datetime
import shutil
import pytesseract
import cv2
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from pdf2image import convert_from_bytes
from PIL import Image
import streamlit as st

# ------------------- Translation backends (try imports) -------------------
translator_backend = None
try:
    # googletrans (unofficial) - requires internet
    from googletrans import Translator as GoogleTranslator
    google_trans = GoogleTranslator()
    translator_backend = "googletrans"
except Exception:
    google_trans = None
    try:
        # deep-translator (fallback)
        from deep_translator import GoogleTranslator as DeepTranslator
        translator_backend = "deep-translator"
    except Exception:
        translator_backend = None

# ------------------- Backward-compatible PDF exception -------------------
try:
    from pdf2image.exceptions import PDFInfoNotInstalledError
except Exception:
    class PDFInfoNotInstalledError(Exception):
        pass

# ------------------- Auto-detect helpers -------------------
def detect_poppler():
    if shutil.which("pdfinfo"):
        return None
    user_dl = os.path.join(os.path.expanduser("~"), "Downloads")
    candidates = [
        r"C:\Program Files\poppler\Library\bin",
        r"C:\Program Files\poppler\bin",
        os.path.join(user_dl, "Release-25.07.0-0", "poppler-25.07.0", "Library", "bin"),
        os.path.join(user_dl, "poppler-25.07.0", "Library", "bin"),
        os.path.join(user_dl, "poppler-24.08.0", "Library", "bin"),
    ]
    for p in candidates:
        if os.path.exists(os.path.join(p, "pdfinfo.exe")):
            return p
    return None

def detect_tesseract():
    if shutil.which("tesseract"):
        return None
    candidates = [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
        os.path.join(os.path.expanduser("~"), "AppData", "Local", "Programs", "Tesseract-OCR", "tesseract.exe"),
    ]
    for c in candidates:
        if os.path.exists(c):
            return c
    return None

poppler_bin = detect_poppler()
tess_path = detect_tesseract()
if poppler_bin:
    os.environ["PATH"] += os.pathsep + poppler_bin
if tess_path:
    pytesseract.pytesseract.tesseract_cmd = tess_path

# ------------------- Configuration -------------------
OUTPUT_DIR = "output_reports"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Supported target languages and mapping for translator library codes
LANGUAGE_CODES = {
    "Kannada": "kn",
    "Hindi": "hi",
    "Marathi": "mr",
    "Tamil": "ta",
    "Telugu": "te",
    "Malayalam": "ml",
    "Gujarati": "gu",
    "Bengali": "bn",
    "Punjabi": "pa",
    "English": "en"
}

# ------------------- Small UI helpers -------------------
def translate_text(text: str, target_lang_code: str) -> str:
    """
    Translate text to target language code.
    Uses googletrans first, then deep-translator as fallback.
    If no backend available, returns original text with note.
    """
    if not text:
        return ""
    # prefer googletrans if available
    try:
        if translator_backend == "googletrans" and google_trans:
            # googletrans auto-detects; large text sometimes breaks — chunk if large
            # We'll translate in blocks of ~4000 chars to be safe
            parts = []
            max_len = 4000
            start = 0
            while start < len(text):
                chunk = text[start:start+max_len]
                res = google_trans.translate(chunk, dest=target_lang_code)
                parts.append(res.text)
                start += max_len
            return "\n".join(parts)
        elif translator_backend == "deep-translator":
            # deep_translator's GoogleTranslator takes source and target
            deep = DeepTranslator(source='auto', target=target_lang_code)
            # deep-translator may have length limits; translate whole text
            return deep.translate(text)
    except Exception as e:
        # If translation fails for large payloads — attempt smaller chunks with fallback
        try:
            # try chunked googletrans if available
            if translator_backend == "googletrans" and google_trans:
                parts = []
                max_len = 3000
                start = 0
                while start < len(text):
                    chunk = text[start:start+max_len]
                    res = google_trans.translate(chunk, dest=target_lang_code)
                    parts.append(res.text)
                    start += max_len
                return "\n".join(parts)
        except Exception:
            pass
    # No translator available or all attempts failed
    return f"[Translation not available — showing original]\n\n{text}"

# ------------------- OCR / PDF helpers -------------------
def preprocess(pil):
    img = np.array(pil.convert("RGB"))
    gray = cv2.cvtColor(img, cv2.COLOR_RGB2GRAY)
    gray = cv2.medianBlur(gray, 3)
    _, th = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    return Image.fromarray(th)

def pdf_to_images_safe(pdf_bytes):
    try:
        images = convert_from_bytes(pdf_bytes, dpi=200)
        return images, None
    except PDFInfoNotInstalledError as e:
        return None, ("PDF_CONVERT_ERROR", str(e))
    except Exception as e:
        return None, ("PDF_OTHER_ERROR", str(e))

def ocr_text_from_image(pil):
    try:
        pre = preprocess(pil)
        txt = pytesseract.image_to_string(pre, lang="eng")
        return txt
    except Exception as e:
        return f"[OCR ERROR: {e}]"

# ------------------- Hospital report parsing helpers -------------------
def parse_cbc_table(text: str) -> pd.DataFrame:
    """
    Try to parse common CBC rows (Hemoglobin, WBC, Platelets etc.) into a DataFrame.
    This is heuristic-based and may not capture every format.
    """
    rows = []
    # normalize lines
    for line in text.splitlines():
        ln = line.strip()
        if not ln:
            continue
        # basic pattern: Name value (unit) reference
        # we look for known keywords
        for key in ["hemoglobin", "hb", "hematocrit", "hct", "red blood", "rbc", "wbc", "platelet", "platelets", "neutrophils", "lymphocytes", "monocytes", "eosinophils", "basophils", "rdw", "mcv", "mch", "mchc", "glucose", "cholesterol", "creatinine", "bun", "alt", "ast"]:
            if key in ln.lower():
                # find first number in line
                m = re.search(r"(-?\d{1,3}(?:[.,]\d+)?(?:e[+-]?\d+)?)", ln)
                if m:
                    val = m.group(1).replace(",", ".")
                    # attempt to find unit (mg/dL, %, 10^9/L etc.)
                    unit_m = re.search(r"(mg/dl|g/dl|%|10\^9/l|10\^12/l|10\^9/L|10\^12/L|u/l|u/L|pg|fl|g/l)", ln, re.I)
                    unit = unit_m.group(0) if unit_m else ""
                    rows.append((ln.split()[0], val, unit, ln))
                else:
                    # no numeric found - include line as-is
                    rows.append((ln.split()[0], "", "", ln))
                break
    if not rows:
        return pd.DataFrame()
    df = pd.DataFrame(rows, columns=["Token", "Value", "Unit", "FullLine"])
    return df

# ------------------- Streamlit UI -------------------
st.set_page_config(page_title="MediScan.AI v4 — Multi-language Hospital Reports", layout="wide")

# Header
st.markdown("<h1 style='text-align:center;color:#004c99'>🩺 MediScan.AI v4 — Hospital Report Translation</h1>", unsafe_allow_html=True)
st.write("Upload hospital report (PDF or image). The app will extract full text, parse CBC if present, and produce translations into multiple Indian languages.")

# Environment status
with st.expander("Environment & Translation Backend Status", expanded=True):
    st.write("**Tesseract:**", "Found ✅" if (shutil.which("tesseract") or tess_path) else "Not found ❌")
    st.write("**Poppler:**", "Found ✅" if (shutil.which("pdfinfo") or poppler_bin) else "Not found ❌")
    st.write("**Translation backend:**", translator_backend if translator_backend else "None (install googletrans or deep-translator for translations)")

upload = st.file_uploader("Upload hospital report (PDF / JPG / PNG)", type=["pdf", "png", "jpg", "jpeg"])

if upload:
    # handle PDF or image
    images = []
    if upload.type == "application/pdf" or upload.name.lower().endswith(".pdf"):
        images, err = pdf_to_images_safe(upload.read())
        if err:
            st.error("Could not convert PDF to images. Install Poppler or upload report as image.")
            st.stop()
    else:
        try:
            images = [Image.open(upload).convert("RGB")]
        except Exception as e:
            st.error(f"Could not open uploaded image: {e}")
            st.stop()

    # OCR entire report
    full_text = ""
    for i, img in enumerate(images):
        with st.spinner(f"OCR page {i+1}..."):
            page_text = ocr_text_from_image(img)
            full_text += page_text + "\n"
            st.markdown(f"**Page {i+1} OCR Preview:**")
            st.text_area(f"OCR Page {i+1}", value=page_text, height=180)

    # show raw extracted text
    st.markdown("### Full Extracted Report Text")
    st.text_area("Full Report Text", value=full_text, height=300)

    # parse CBC / lab table heuristically
    df_cbc = parse_cbc_table(full_text)
    if not df_cbc.empty:
        st.markdown("### Parsed CBC / Lab-like rows (heuristic parse)")
        st.dataframe(df_cbc)
        # save parsed CSV
        ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        parsed_csv_path = os.path.join(OUTPUT_DIR, f"report_{ts}_parsed_lab.csv")
        df_cbc.to_csv(parsed_csv_path, index=False)
        st.success(f"Parsed table saved: {parsed_csv_path}")
    else:
        st.info("No structured CBC rows detected automatically. Full text is available above for manual review.")

    # Translation controls
    st.markdown("---")
    st.markdown("## Translate Report")
    languages = list(LANGUAGE_CODES.keys())
    languages.remove("English")  # English is original
    # include multi-select and All button
    choose = st.multiselect("Select target languages (Ctrl/Command-click to multi-select):", languages, default=["Hindi"])
    if st.checkbox("Translate to ALL supported languages"):
        choose = languages.copy()

    # Translate and save
    if st.button("Generate Translations"):
        if not choose:
            st.warning("Choose at least one language.")
        else:
            ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            saved_files = []
            for lang in choose:
                code = LANGUAGE_CODES.get(lang, "en")
                st.markdown(f"Translating to **{lang}** ...")
                translated = translate_text(full_text, code)
                # show translated in expander
                with st.expander(f"{lang} translation preview"):
                    st.text_area(f"{lang} translation", value=translated, height=300)
                # save to file
                fname = os.path.join(OUTPUT_DIR, f"report_{ts}_translated_{code}.txt")
                with open(fname, "w", encoding="utf-8") as f:
                    f.write(f"--- ORIGINAL REPORT (extracted) ---\n\n{full_text}\n\n")
                    f.write(f"--- TRANSLATION: {lang} ({code}) ---\n\n{translated}\n")
                saved_files.append((lang, code, fname))
                st.success(f"Saved translation: {fname}")

            # Offer downloads for saved files
            st.markdown("### Download Translated Files")
            for lang, code, path in saved_files:
                with open(path, "rb") as f:
                    data = f.read()
                st.download_button(label=f"Download {lang} (.txt)", data=data, file_name=os.path.basename(path), mime="text/plain")

    # Option: create a combined PPT slide with the original + one translation screenshot placeholder
    if st.button("Create quick PPT with original + first translation placeholder"):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            prs = Presentation()
            # slide 1: title
            s0 = prs.slides.add_slide(prs.slide_layouts[0])
            s0.shapes.title.text = "MediScan.AI — Report Translation"
            s0.placeholders[1].text = "Auto-generated report snapshot"
            # slide 2: original text
            s1 = prs.slides.add_slide(prs.slide_layouts[1])
            s1.shapes.title.text = "Original Extracted Text"
            tf = s1.shapes.placeholders[1].text_frame
            # add first 15 lines to avoid overflow; user can edit later
            lines = full_text.splitlines()
            for ln in lines[:30]:
                p = tf.add_paragraph()
                p.text = ln[:100]
                p.font.size = Pt(12)
            # save
            ppt_path = os.path.join(OUTPUT_DIR, f"report_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}_quick.pptx")
            prs.save(ppt_path)
            with open(ppt_path, "rb") as f:
                st.download_button("Download PPTX", data=f.read(), file_name=os.path.basename(ppt_path))
            st.success(f"PPT created: {ppt_path}")
        except Exception as e:
            st.error(f"Could not create PPT automatically: {e}")

st.markdown("---")
st.caption("© MediScan.AI v4 — Team DiamondHot | Note: Translations use online services if installed. Always consult a clinician for medical decisions.")
